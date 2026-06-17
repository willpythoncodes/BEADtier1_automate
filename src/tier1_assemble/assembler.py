from __future__ import annotations

import json
import os
import re
import tempfile
from dataclasses import replace
from pathlib import Path
from typing import Dict, Optional

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .geometry import Box, contain_box
from .models import (
    AssetBinding,
    AssemblyManifest,
    AssemblyOptions,
    AssemblyReport,
    TokenReplacement,
)


TOKEN_RE = re.compile(r"\{([A-Za-z0-9_]+)\}")
IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg"}
DEFAULT_ASSET_DIR_NAME = "placeholder_assets"
DEFAULT_ASSET_DIR_RUN_ID = "placeholder_review"
COMPOSITE_POSITIONS = {"bottom-right", "bottom-left", "top-right", "top-left"}
COMPOSITE_MARGIN_PX = 50


def load_assembly_manifest(
    manifest_path: Path,
    scratch_dir: Optional[Path] = None,
    run_id: Optional[str] = None,
) -> AssemblyManifest:
    manifest_path = Path(manifest_path)
    if not manifest_path.exists():
        raise FileNotFoundError(f"assembly manifest does not exist: {manifest_path}")
    if manifest_path.suffix.lower() != ".json":
        raise ValueError("assembly manifest path must have a .json extension")

    try:
        with manifest_path.open("r", encoding="utf-8") as handle:
            payload = json.load(handle)
    except json.JSONDecodeError as exc:
        raise ValueError(f"invalid JSON manifest: {exc}") from exc

    if not isinstance(payload, dict):
        raise ValueError("assembly manifest must be a JSON object")
    version = payload.get("version", 1)
    if version != 1:
        raise ValueError(f"unsupported assembly manifest version: {version}")

    manifest_run_id = run_id or payload.get("run_id")
    if not isinstance(manifest_run_id, str) or not manifest_run_id.strip():
        raise ValueError("assembly manifest requires a non-empty run_id")
    manifest_run_id = manifest_run_id.strip()

    assets_payload = payload.get("assets")
    if not isinstance(assets_payload, dict) or not assets_payload:
        raise ValueError("assembly manifest must contain a non-empty assets object")

    assets_dir = Path(scratch_dir).expanduser().resolve() if scratch_dir else manifest_path.resolve().parent
    assets: Dict[str, AssetBinding] = {}
    for raw_token, asset_payload in assets_payload.items():
        token = _normalize_token(raw_token)
        if not token:
            raise ValueError("asset token keys must be non-empty strings")
        if token in assets:
            raise ValueError(f"duplicate asset token: {token}")
        assets[token] = _parse_asset_binding(token, asset_payload, assets_dir)

    return AssemblyManifest(
        manifest_path=manifest_path.resolve(),
        assets_dir=assets_dir,
        run_id=manifest_run_id,
        assets=assets,
    )


def load_asset_directory(asset_dir: Path, run_id: Optional[str] = None) -> AssemblyManifest:
    asset_dir = Path(asset_dir).expanduser().resolve()
    if not asset_dir.exists():
        raise FileNotFoundError(f"asset directory does not exist: {asset_dir}")
    if not asset_dir.is_dir():
        raise ValueError(f"asset directory path must be a directory: {asset_dir}")

    assets: Dict[str, AssetBinding] = {}
    for path in sorted(asset_dir.iterdir()):
        if not path.is_file() or path.suffix.lower() not in IMAGE_SUFFIXES:
            continue
        token = _normalize_token(path.stem)
        if not token:
            continue
        if token in assets:
            raise ValueError(f"duplicate asset token inferred from filenames: {token}")
        assets[token] = AssetBinding(token=token, path=path.resolve(), kind="image")

    if not assets:
        raise ValueError(f"asset directory contains no PNG or JPEG files: {asset_dir}")

    manifest_run_id = (run_id or DEFAULT_ASSET_DIR_RUN_ID).strip()
    if not manifest_run_id:
        raise ValueError("asset directory assembly requires a non-empty run_id")

    return AssemblyManifest(
        manifest_path=asset_dir,
        assets_dir=asset_dir,
        run_id=manifest_run_id,
        assets=assets,
    )


def assemble_presentation(
    template_path: Path,
    manifest_path: Optional[Path],
    output_dir: Path,
    scratch_dir: Optional[Path] = None,
    asset_dir: Optional[Path] = None,
    run_id: Optional[str] = None,
    output_path: Optional[Path] = None,
    options: Optional[AssemblyOptions] = None,
) -> AssemblyReport:
    options = options or AssemblyOptions()
    template_path = Path(template_path)
    output_dir = Path(output_dir)
    manifest = _load_asset_source(template_path, manifest_path, scratch_dir, asset_dir, run_id)
    if asset_dir:
        manifest = replace(manifest, assets_dir=Path(asset_dir).expanduser().resolve())
    output_path = Path(output_path) if output_path else output_dir / f"{manifest.run_id}_Tier1_Report_Draft.pptx"
    _guard_assembly_paths(template_path, output_path, options)

    presentation = Presentation(str(template_path))
    report = AssemblyReport(
        template_path=template_path.resolve(),
        manifest_path=manifest.manifest_path,
        output_path=output_path.resolve(),
        run_id=manifest.run_id,
    )

    used_assets = set()
    for slide_number, slide in enumerate(presentation.slides, start=1):
        report.slides_scanned += 1
        for container, shape, tokens, group_path in _iter_token_shapes(slide.shapes, slide):
            report.tokens_found += len(tokens)
            shape_label = _shape_label(shape, group_path)
            if len(tokens) > 1:
                report.warnings.append(
                    f"slide {slide_number} shape {shape_label} has multiple tokens; using {tokens[0]}"
                )
            token = tokens[0]
            replacement = TokenReplacement(
                token=token,
                slide_number=slide_number,
                shape_name=shape_label,
                visible_text_token=_shape_visible_text_contains(shape, token),
            )
            report.replacements.append(replacement)

            binding = manifest.assets.get(token)
            if not binding:
                report.unresolved_tokens.append(token)
                replacement.warning = "no matching asset binding"
                continue

            try:
                asset_path = _prepare_asset(binding, manifest, report)
                _replace_shape_with_picture(container, shape, asset_path, options)
            except Exception as exc:
                warning = f"token {token} on slide {slide_number}: {exc}"
                report.warnings.append(warning)
                replacement.warning = str(exc)
                continue

            used_assets.add(token)
            replacement.asset_path = asset_path
            replacement.inserted = True
            report.assets_inserted += 1

    report.unresolved_tokens = sorted(set(report.unresolved_tokens))
    report.unused_assets = sorted(set(manifest.assets) - used_assets)
    _save_atomically(presentation, output_path)
    return report


def combine_map_legend(
    map_path: Path,
    legend_path: Path,
    position: str,
    output_path: Optional[Path] = None,
) -> Path:
    map_path = Path(map_path).expanduser().resolve()
    legend_path = Path(legend_path).expanduser().resolve()
    position = _normalize_composite_anchor(position)
    _validate_image_asset(map_path)
    _validate_image_asset(legend_path)

    output_path = Path(output_path).expanduser().resolve() if output_path else _default_composite_output_path(map_path)
    if output_path.suffix.lower() != ".png":
        raise ValueError("combined output path must have a .png extension")
    output_path.parent.mkdir(parents=True, exist_ok=True)

    base = Image.open(map_path).convert("RGBA")
    overlay = Image.open(legend_path).convert("RGBA")
    result = _composite_images(base, overlay, position, COMPOSITE_MARGIN_PX)
    result.save(output_path)
    return output_path


def _load_asset_source(
    template_path: Path,
    manifest_path: Optional[Path],
    scratch_dir: Optional[Path],
    asset_dir: Optional[Path],
    run_id: Optional[str],
) -> AssemblyManifest:
    if manifest_path:
        return load_assembly_manifest(Path(manifest_path), scratch_dir=scratch_dir, run_id=run_id)

    if asset_dir:
        return load_asset_directory(Path(asset_dir), run_id=run_id)

    if scratch_dir:
        return load_asset_directory(Path(scratch_dir) / DEFAULT_ASSET_DIR_NAME, run_id=run_id)

    return load_asset_directory(template_path.resolve().parent / DEFAULT_ASSET_DIR_NAME, run_id=run_id)


def write_assembly_report(report: AssemblyReport, report_path: Path) -> None:
    report_path = Path(report_path)
    if report_path.suffix.lower() != ".json":
        raise ValueError("assembly report path must have a .json extension")
    report_path.parent.mkdir(parents=True, exist_ok=True)
    temp_file = tempfile.NamedTemporaryFile(
        prefix=f".{report_path.name}.",
        suffix=".tmp",
        dir=str(report_path.parent),
        delete=False,
        mode="w",
        encoding="utf-8",
    )
    temp_path = Path(temp_file.name)
    try:
        with temp_file:
            json.dump(report.as_dict(), temp_file, indent=2, sort_keys=True)
            temp_file.write("\n")
        os.replace(str(temp_path), str(report_path))
    except Exception:
        temp_path.unlink(missing_ok=True)
        raise


def _parse_asset_binding(token: str, payload, asset_root: Path) -> AssetBinding:
    if isinstance(payload, str):
        path = _resolve_asset_path(asset_root, payload)
        return AssetBinding(token=token, path=path, kind=_infer_asset_kind(path))

    if not isinstance(payload, dict):
        raise ValueError(f"asset {token} must be a string path or object")

    kind = payload.get("type", "image")
    if kind in {"image", "html"}:
        path_value = payload.get("path")
        if not isinstance(path_value, str) or not path_value.strip():
            raise ValueError(f"asset {token} requires a non-empty path")
        return AssetBinding(
            token=token,
            path=_resolve_asset_path(asset_root, path_value),
            kind=kind,
        )

    if kind == "composite":
        base_value = payload.get("base")
        overlay_value = payload.get("overlay")
        if not isinstance(base_value, str) or not base_value.strip():
            raise ValueError(f"composite asset {token} requires a non-empty base")
        if not isinstance(overlay_value, str) or not overlay_value.strip():
            raise ValueError(f"composite asset {token} requires a non-empty overlay")
        anchor = _normalize_composite_anchor(str(payload.get("anchor", "bottom-right")), token=token)
        return AssetBinding(
            token=token,
            kind="composite",
            base_path=_resolve_asset_path(asset_root, base_value),
            overlay_path=_resolve_asset_path(asset_root, overlay_value),
            anchor=anchor,
            margin_px=int(payload.get("margin_px", 50)),
        )

    raise ValueError(f"unsupported asset type for {token}: {kind}")


def _prepare_asset(binding: AssetBinding, manifest: AssemblyManifest, report: AssemblyReport) -> Path:
    if binding.kind == "composite":
        return _prepare_composite_asset(binding, manifest, report)
    if binding.kind == "html" or (binding.path and binding.path.suffix.lower() in {".html", ".htm"}):
        raise ValueError("HTML table assets must be pre-rendered to PNG for this local v1")
    if not binding.path:
        raise ValueError("asset binding has no path")
    _validate_image_asset(binding.path)
    return binding.path


def _prepare_composite_asset(binding: AssetBinding, manifest: AssemblyManifest, report: AssemblyReport) -> Path:
    if not binding.base_path or not binding.overlay_path:
        raise ValueError("composite asset requires base and overlay paths")
    _validate_image_asset(binding.base_path)
    _validate_image_asset(binding.overlay_path)

    manifest.assets_dir.mkdir(parents=True, exist_ok=True)
    output_path = manifest.assets_dir / f"{binding.token}.png"

    base = Image.open(binding.base_path).convert("RGBA")
    overlay = Image.open(binding.overlay_path).convert("RGBA")
    result = _composite_images(base, overlay, binding.anchor, binding.margin_px)
    result.save(output_path)
    report.derived_assets_created += 1
    return output_path


def _composite_images(base: Image.Image, overlay: Image.Image, anchor: str, margin_px: int) -> Image.Image:
    bbox = overlay.getchannel("A").getbbox()
    if bbox is None:
        raise ValueError("overlay has no visible pixels")
    overlay_crop = overlay.crop(bbox)
    x, y = _anchored_overlay_position(base.size, overlay_crop.size, anchor, margin_px)
    result = base.copy()
    result.alpha_composite(overlay_crop, dest=(x, y))
    return result


def _anchored_overlay_position(base_size: tuple, overlay_size: tuple, anchor: str, margin: int) -> tuple:
    base_w, base_h = base_size
    overlay_w, overlay_h = overlay_size
    margin = max(0, margin)

    if anchor in {"top-left", "left", "bottom-left"}:
        x = margin
    elif anchor in {"top-right", "right", "bottom-right"}:
        x = base_w - overlay_w - margin
    else:
        x = (base_w - overlay_w) // 2

    if anchor in {"top-left", "top", "top-right"}:
        y = margin
    elif anchor in {"bottom-left", "bottom", "bottom-right"}:
        y = base_h - overlay_h - margin
    else:
        y = (base_h - overlay_h) // 2

    if x < 0 or y < 0:
        raise ValueError("overlay does not fit inside base image with requested margin")
    return x, y


def _iter_token_shapes(shapes, container, group_path: tuple = ()):
    for shape in list(shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            group_name = getattr(shape, "name", "<unnamed>")
            yield from _iter_token_shapes(shape.shapes, shape, group_path + (group_name,))
            continue
        tokens = _shape_tokens(shape)
        if tokens:
            yield container, shape, tokens, group_path


def _shape_label(shape, group_path: tuple) -> str:
    shape_name = getattr(shape, "name", "<unnamed>")
    if not group_path:
        return shape_name
    return "/".join(group_path + (shape_name,))


def _replace_shape_with_picture(container, shape, asset_path: Path, options: AssemblyOptions) -> None:
    with Image.open(asset_path) as image:
        image_size = image.size
    slot = Box(shape.left, shape.top, shape.width, shape.height)
    box = contain_box(slot, image_size, options.anchor, options.margin)
    shape._element.getparent().remove(shape._element)
    container.shapes.add_picture(str(asset_path), box.left, box.top, width=box.width, height=box.height)


def _shape_tokens(shape) -> list:
    values = []
    c_nv_prs = shape._element.xpath(".//p:cNvPr")
    if c_nv_prs:
        values.extend(filter(None, [c_nv_prs[0].get("title"), c_nv_prs[0].get("descr")]))
    if getattr(shape, "has_text_frame", False):
        values.append(shape.text)
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            values.extend(cell.text for cell in row.cells)

    tokens = []
    for value in values:
        for match in TOKEN_RE.finditer(value or ""):
            token = match.group(1)
            if token not in tokens:
                tokens.append(token)
    return tokens


def _shape_visible_text_contains(shape, token: str) -> bool:
    wrapped = f"{{{token}}}"
    if getattr(shape, "has_text_frame", False) and wrapped in shape.text:
        return True
    if getattr(shape, "has_table", False):
        return any(wrapped in cell.text for row in shape.table.rows for cell in row.cells)
    return False


def _guard_assembly_paths(template_path: Path, output_path: Path, options: AssemblyOptions) -> None:
    if not template_path.exists():
        raise FileNotFoundError(f"template PPTX does not exist: {template_path}")
    if template_path.suffix.lower() != ".pptx":
        raise ValueError("template path must have a .pptx extension")
    if output_path.suffix.lower() != ".pptx":
        raise ValueError("output path must have a .pptx extension")
    if template_path.resolve() == output_path.resolve():
        raise ValueError("output path must be different from template path")
    if output_path.exists() and not options.overwrite:
        raise FileExistsError(f"output already exists: {output_path}")


def _validate_image_asset(path: Path) -> None:
    if not path.exists():
        raise FileNotFoundError(f"asset does not exist: {path}")
    if path.suffix.lower() not in IMAGE_SUFFIXES:
        raise ValueError(f"asset must be a PNG or JPEG image: {path}")


def _infer_asset_kind(path: Path) -> str:
    if path.suffix.lower() in {".html", ".htm"}:
        return "html"
    return "image"


def _resolve_asset_path(root: Path, value: str) -> Path:
    path = Path(value).expanduser()
    if path.is_absolute():
        return path.resolve()
    return (root / path).resolve()


def _normalize_composite_anchor(anchor: str, token: Optional[str] = None) -> str:
    anchor = anchor.strip().lower()
    if anchor not in COMPOSITE_POSITIONS:
        valid = ", ".join(sorted(COMPOSITE_POSITIONS))
        if token:
            raise ValueError(f"composite asset {token}: position must be one of: {valid}")
        raise ValueError(f"position must be one of: {valid}")
    return anchor


def _default_composite_output_path(map_path: Path) -> Path:
    return map_path.with_name(f"{map_path.stem}_with_legend.png")


def _normalize_token(value: str) -> str:
    if not isinstance(value, str):
        return ""
    value = value.strip()
    if value.startswith("{") and value.endswith("}"):
        value = value[1:-1]
    return value.strip()


def _save_atomically(presentation, output_path: Path) -> None:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    temp_file = tempfile.NamedTemporaryFile(
        prefix=f".{output_path.name}.",
        suffix=".tmp",
        dir=str(output_path.parent),
        delete=False,
    )
    temp_path = Path(temp_file.name)
    temp_file.close()
    try:
        presentation.save(str(temp_path))
        os.replace(str(temp_path), str(output_path))
    except Exception:
        temp_path.unlink(missing_ok=True)
        raise
