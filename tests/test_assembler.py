import json
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from tier1_assemble.assembler import combine_map_legend, assemble_presentation, load_assembly_manifest, load_asset_directory


def test_load_assembly_manifest_resolves_assets_relative_to_scratch_dir(tmp_path):
    scratch = tmp_path / "scratch" / "run-a"
    scratch.mkdir(parents=True)
    image = scratch / "map.png"
    _make_image(image)
    manifest = tmp_path / "assets.json"
    _write_manifest(manifest, {"run_id": "run-a", "assets": {"hazard_map": "map.png"}})

    parsed = load_assembly_manifest(manifest, scratch_dir=scratch)

    assert parsed.run_id == "run-a"
    assert parsed.assets_dir == scratch.resolve()
    assert parsed.assets["hazard_map"].path == image.resolve()


def test_load_asset_directory_infers_tokens_from_file_stems(tmp_path):
    asset_dir = tmp_path / "placeholder_assets"
    asset_dir.mkdir()
    image = asset_dir / "hazard_map.png"
    ignored = asset_dir / ".DS_Store"
    _make_image(image)
    ignored.write_text("ignored", encoding="utf-8")

    parsed = load_asset_directory(asset_dir, run_id="run-a")

    assert parsed.run_id == "run-a"
    assert parsed.manifest_path == asset_dir.resolve()
    assert parsed.assets_dir == asset_dir.resolve()
    assert parsed.assets["hazard_map"].path == image.resolve()
    assert ".DS_Store" not in parsed.assets


def test_assemble_replaces_token_shapes_and_preserves_narrative_text(tmp_path):
    template = tmp_path / "template.pptx"
    image = tmp_path / "map.png"
    table_image = tmp_path / "table.png"
    manifest = tmp_path / "assets.json"
    output_dir = tmp_path / "output"
    _make_image(image, size=(800, 400))
    _make_image(table_image, size=(300, 300))
    _make_template(template)
    _write_manifest(
        manifest,
        {
            "run_id": "run-a",
            "assets": {
                "hazard_map": "map.png",
                "summary_table": "table.png",
            },
        },
    )

    report = assemble_presentation(template, manifest, output_dir)

    assert report.valid
    assert report.output_path == (output_dir / "run-a_Tier1_Report_Draft.pptx").resolve()
    assert report.tokens_found == 2
    assert report.assets_inserted == 2
    assert report.unused_assets == []
    reopened = Presentation(str(report.output_path))
    slide = reopened.slides[0]
    pictures = [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pictures) == 2
    assert any(shape.has_text_frame and shape.text == "Manual narrative remains editable" for shape in slide.shapes)
    assert all("{hazard_map}" not in getattr(shape, "text", "") for shape in slide.shapes if getattr(shape, "has_text_frame", False))


def test_assemble_can_use_placeholder_assets_directory_without_manifest(tmp_path):
    template = tmp_path / "template.pptx"
    scratch = tmp_path / "scratch"
    asset_dir = scratch / "placeholder_assets"
    output_dir = tmp_path / "output"
    asset_dir.mkdir(parents=True)
    _make_image(asset_dir / "hazard_map.png", size=(800, 400))
    _make_image(asset_dir / "summary_table.png", size=(300, 300))
    _make_template(template)

    report = assemble_presentation(template, None, output_dir, scratch_dir=scratch, run_id="run-a")

    assert report.valid
    assert report.manifest_path == asset_dir.resolve()
    assert report.output_path == (output_dir / "run-a_Tier1_Report_Draft.pptx").resolve()
    assert report.assets_inserted == 2


def test_load_assembly_manifest_rejects_invalid_composite_anchor(tmp_path):
    manifest = tmp_path / "assets.json"
    _write_manifest(
        manifest,
        {
            "run_id": "run-a",
            "assets": {
                "map_with_legend": {
                    "type": "composite",
                    "base": "base.png",
                    "overlay": "overlay.png",
                    "anchor": "botom-right",
                }
            },
        },
    )

    with pytest.raises(ValueError, match=r"composite asset map_with_legend: position must be one of: bottom-left, bottom-right, top-left, top-right"):
        load_assembly_manifest(manifest)


def test_combine_map_legend_rejects_invalid_anchor(tmp_path):
    map_path = tmp_path / "map.png"
    legend_path = tmp_path / "legend.png"
    _make_image(map_path)
    _make_image(legend_path)

    with pytest.raises(ValueError, match=r"position must be one of: bottom-left, bottom-right, top-left, top-right"):
        combine_map_legend(map_path, legend_path, "botom-right")


def test_assemble_replaces_token_shapes_inside_groups(tmp_path):
    template = tmp_path / "template.pptx"
    image = tmp_path / "map.png"
    manifest = tmp_path / "assets.json"
    output_dir = tmp_path / "output"
    _make_image(image, size=(800, 400))
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    group = slide.shapes.add_group_shape()
    group.name = "Map Frame"
    token_box = group.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(4), Inches(2))
    token_box.name = "Map Placeholder"
    token_box.text = "{hazard_map}"
    _set_alt_text(token_box, "{hazard_map}")
    presentation.save(str(template))
    _write_manifest(manifest, {"run_id": "run-a", "assets": {"hazard_map": "map.png"}})

    report = assemble_presentation(template, manifest, output_dir)

    assert report.valid
    assert report.assets_inserted == 1
    assert report.replacements[0].shape_name == "Map Frame/Map Placeholder"
    reopened = Presentation(str(report.output_path))
    grouped = [shape for shape in reopened.slides[0].shapes if shape.shape_type == MSO_SHAPE_TYPE.GROUP]
    assert len(grouped) == 1
    pictures = [shape for shape in grouped[0].shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pictures) == 1


def test_assemble_replaces_token_shapes_inside_nested_groups(tmp_path):
    template = tmp_path / "template.pptx"
    image = tmp_path / "map.png"
    manifest = tmp_path / "assets.json"
    output_dir = tmp_path / "output"
    _make_image(image, size=(800, 400))
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    outer = slide.shapes.add_group_shape()
    outer.name = "Outer Group"
    inner = outer.shapes.add_group_shape()
    inner.name = "Inner Group"
    token_box = inner.shapes.add_textbox(Inches(0.25), Inches(0.25), Inches(3), Inches(1.5))
    token_box.name = "Map Placeholder"
    token_box.text = "{hazard_map}"
    _set_alt_text(token_box, "{hazard_map}")
    presentation.save(str(template))
    _write_manifest(manifest, {"run_id": "run-a", "assets": {"hazard_map": "map.png"}})

    report = assemble_presentation(template, manifest, output_dir)

    assert report.valid
    assert report.replacements[0].shape_name == "Outer Group/Inner Group/Map Placeholder"
    reopened = Presentation(str(report.output_path))
    outer_group = next(shape for shape in reopened.slides[0].shapes if shape.shape_type == MSO_SHAPE_TYPE.GROUP)
    inner_group = next(shape for shape in outer_group.shapes if shape.shape_type == MSO_SHAPE_TYPE.GROUP)
    pictures = [shape for shape in inner_group.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pictures) == 1


def test_assemble_valid_despite_informational_warnings(tmp_path):
    template = tmp_path / "template.pptx"
    image = tmp_path / "map.png"
    manifest = tmp_path / "assets.json"
    output_dir = tmp_path / "output"
    _make_image(image, size=(800, 400))
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    token_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
    token_box.text = "{hazard_map}{summary_table}"
    _set_alt_text(token_box, "{hazard_map}{summary_table}")
    presentation.save(str(template))
    _write_manifest(manifest, {"run_id": "run-a", "assets": {"hazard_map": "map.png"}})

    report = assemble_presentation(template, manifest, output_dir)

    assert report.valid
    assert report.assets_inserted == 1
    assert any("multiple tokens" in warning for warning in report.warnings)


def test_assemble_dedupes_unresolved_tokens(tmp_path):
    template = tmp_path / "template.pptx"
    manifest = tmp_path / "assets.json"
    output_dir = tmp_path / "output"
    presentation = Presentation()
    for _ in range(2):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        token_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
        token_box.text = "{hazard_map}"
        _set_alt_text(token_box, "{hazard_map}")
    presentation.save(str(template))
    _write_manifest(manifest, {"run_id": "run-a", "assets": {"other_token": "missing.png"}})

    report = assemble_presentation(template, manifest, output_dir)

    assert not report.valid
    assert report.unresolved_tokens == ["hazard_map"]


def test_assemble_reports_unresolved_tokens(tmp_path):
    template = tmp_path / "template.pptx"
    manifest = tmp_path / "assets.json"
    output_dir = tmp_path / "output"
    _make_template(template)
    _write_manifest(manifest, {"run_id": "run-a", "assets": {"other_token": "missing.png"}})

    report = assemble_presentation(template, manifest, output_dir)

    assert not report.valid
    assert "hazard_map" in report.unresolved_tokens
    assert "summary_table" in report.unresolved_tokens


def test_assemble_creates_default_composite_asset(tmp_path):
    template = tmp_path / "template.pptx"
    base = tmp_path / "base.png"
    overlay = tmp_path / "overlay.png"
    manifest = tmp_path / "assets.json"
    output_dir = tmp_path / "output"
    _make_single_token_template(template, "{map_with_legend}")
    Image.new("RGBA", (300, 200), (255, 255, 255, 255)).save(base)
    legend = Image.new("RGBA", (300, 200), (0, 0, 0, 0))
    for x in range(220, 280):
        for y in range(20, 60):
            legend.putpixel((x, y), (255, 0, 0, 255))
    legend.save(overlay)
    _write_manifest(
        manifest,
        {
            "run_id": "run-a",
            "assets": {
                "map_with_legend": {
                    "type": "composite",
                    "base": "base.png",
                    "overlay": "overlay.png",
                }
            },
        },
    )

    report = assemble_presentation(template, manifest, output_dir)

    assert report.valid
    assert report.derived_assets_created == 1
    composite = tmp_path / "map_with_legend.png"
    assert composite.exists()
    result = Image.open(composite).convert("RGBA")
    assert result.getpixel((190, 110)) == (255, 0, 0, 255)


def test_assemble_writes_composite_assets_to_asset_dir(tmp_path):
    template = tmp_path / "template.pptx"
    asset_dir = tmp_path / "placeholder_assets"
    output_dir = tmp_path / "output"
    asset_dir.mkdir()
    base = asset_dir / "base.png"
    overlay = asset_dir / "overlay.png"
    _make_single_token_template(template, "{map_with_legend}")
    Image.new("RGBA", (300, 200), (255, 255, 255, 255)).save(base)
    legend = Image.new("RGBA", (300, 200), (0, 0, 0, 0))
    for x in range(220, 280):
        for y in range(20, 60):
            legend.putpixel((x, y), (255, 0, 0, 255))
    legend.save(overlay)
    manifest = tmp_path / "assets.json"
    _write_manifest(
        manifest,
        {
            "run_id": "run-a",
            "assets": {
                "map_with_legend": {
                    "type": "composite",
                    "base": "placeholder_assets/base.png",
                    "overlay": "placeholder_assets/overlay.png",
                }
            },
        },
    )

    report = assemble_presentation(template, manifest, output_dir, asset_dir=asset_dir)

    assert report.valid
    composite = asset_dir / "map_with_legend.png"
    assert composite.exists()


def test_assemble_rejects_html_assets_until_prerendered(tmp_path):
    template = tmp_path / "template.pptx"
    html = tmp_path / "table.html"
    manifest = tmp_path / "assets.json"
    output_dir = tmp_path / "output"
    _make_single_token_template(template, "{summary_table}")
    html.write_text("<table><tr><td>A</td></tr></table>", encoding="utf-8")
    _write_manifest(manifest, {"run_id": "run-a", "assets": {"summary_table": "table.html"}})

    report = assemble_presentation(template, manifest, output_dir)

    assert not report.valid
    assert "HTML table assets must be pre-rendered" in report.warnings[0]


def test_combine_map_legend_uses_requested_corner_and_50px_margin(tmp_path):
    map_path = tmp_path / "map.png"
    legend_path = tmp_path / "legend.png"
    output_path = tmp_path / "combined.png"
    Image.new("RGBA", (300, 200), (255, 255, 255, 255)).save(map_path)
    legend = Image.new("RGBA", (120, 80), (0, 0, 0, 0))
    for x in range(20, 50):
        for y in range(10, 30):
            legend.putpixel((x, y), (255, 0, 0, 255))
    legend.save(legend_path)

    result_path = combine_map_legend(map_path, legend_path, "top-left", output_path=output_path)

    assert result_path == output_path.resolve()
    result = Image.open(result_path).convert("RGBA")
    assert result.getpixel((50, 50)) == (255, 0, 0, 255)
    assert result.getpixel((49, 49)) == (255, 255, 255, 255)


def _make_template(path: Path):
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    token_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
    token_box.name = "Map Placeholder"
    token_box.text = "{hazard_map}"
    _set_alt_text(token_box, "{hazard_map}")
    table = slide.shapes.add_table(2, 2, Inches(6), Inches(1), Inches(3), Inches(2))
    table.name = "Table Placeholder"
    _set_alt_text(table, "{summary_table}")
    narrative = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(5), Inches(1))
    narrative.text = "Manual narrative remains editable"
    presentation.save(str(path))


def _make_single_token_template(path: Path, token: str):
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    token_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(3))
    token_box.text = token
    _set_alt_text(token_box, token)
    presentation.save(str(path))


def _set_alt_text(shape, token: str):
    c_nv_pr = shape._element.xpath(".//p:cNvPr")[0]
    c_nv_pr.set("title", token)
    c_nv_pr.set("descr", token)


def _make_image(path: Path, size=(400, 200)):
    Image.new("RGB", size, color=(20, 100, 180)).save(path)


def _write_manifest(path: Path, payload: dict):
    path.write_text(json.dumps(payload), encoding="utf-8")
